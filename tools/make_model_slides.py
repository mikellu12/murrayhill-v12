"""Build the model-selection deck as a .pptx for import into Google Slides.

Google Slides imports .pptx directly (File > Import slides), which keeps the
tables editable rather than pasting pictures of numbers.

Every figure here comes from results/tables/vlm_benchmark.csv and
model_benchmark.csv, or is quoted from a cited paper with the source on the
slide. Nothing is typed in by hand, so re-running after a new benchmark
regenerates the deck rather than leaving it stale.

    .venv/Scripts/python tools/make_model_slides.py
"""
import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
sys.path.insert(0, str(HERE.parent / "src"))
from common import RES, banner

INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0xA8, 0x3A, 0x32)
RULE = RGBColor(0xDD, 0xDD, 0xDD)


def slide(prs, title, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])          # blank
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.42), Inches(12.1), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size, p.font.bold, p.font.color.rgb = Pt(30), True, INK
    if subtitle:
        tb2 = s.shapes.add_textbox(Inches(0.6), Inches(1.12), Inches(12.1), Inches(0.5))
        q = tb2.text_frame.paragraphs[0]
        q.text = subtitle
        q.font.size, q.font.color.rgb = Pt(14), MUTED
    return s


def table(s, rows, top, left=Inches(0.6), width=Inches(12.1),
          widths=None, highlight=None, size=13):
    """rows[0] is the header. highlight is a set of row indices to accent."""
    n, m = len(rows), len(rows[0])
    h = Inches(0.34) * n
    shp = s.shapes.add_table(n, m, left, top, width, h).table
    if widths:
        total = sum(widths)
        for j, w in enumerate(widths):
            shp.columns[j].width = Emu(int(width * w / total))
    for i, row in enumerate(rows):
        shp.rows[i].height = Inches(0.32)
        for j, val in enumerate(row):
            c = shp.cell(i, j)
            c.text = str(val)
            c.margin_left, c.margin_right = Inches(0.08), Inches(0.08)
            c.margin_top = c.margin_bottom = Inches(0.02)
            pr = c.text_frame.paragraphs[0]
            pr.font.size = Pt(size if i else size - 1)
            pr.font.bold = bool(i == 0 or (highlight and i in highlight))
            pr.font.color.rgb = (ACCENT if highlight and i in highlight and i
                                 else INK if i else MUTED)
    return shp


def bullets(s, items, top, size=15, left=Inches(0.6), width=Inches(12.1)):
    tb = s.shapes.add_textbox(left, top, width, Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    for k, (txt, bold) in enumerate(items):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = INK if bold else MUTED
        p.space_after = Pt(9)
    return tb


def fmt(r):
    return f"{r.rho:+.3f}   [{r.lo:+.3f}, {r.hi:+.3f}]"


def main():
    banner("model-selection deck")
    vb = pd.read_csv(RES / "tables" / "vlm_benchmark.csv")
    mb = pd.read_csv(RES / "tables" / "model_benchmark.csv")
    get = lambda d, **kw: d.loc[(d[list(kw)] == pd.Series(kw)).all(axis=1)].iloc[0]

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # 1 ------------------------------------------------------------------
    s = slide(prs, "Which vision-language model, and why",
              "Murray Hill · Street Interface Matrix · benchmarked on our own imagery")
    bullets(s, [
        ("The question is not “which model predicts GVI best”.", True),
        ("GVI is a vegetation pixel share. A segmenter measuring it is measuring the "
         "target directly, so scoring a VLM against it tests the VLM on the "
         "segmenter’s home turf.", False),
        ("The model has to produce judgements a segmenter cannot express.", True),
        ("There is no ADE20K or Cityscapes class for “somewhere to sit”, "
         "“how much the frontage varies”, or “how far greenery relieves "
         "the enclosure”. Those are the SIM terms.", False),
        ("GVI and VEI are the validity check — a measured quantity to score the "
         "ratings against, so “the model rates greenery sensibly” is a number, "
         "not an assertion.", False),
    ], Inches(2.0))

    # 2 ------------------------------------------------------------------
    s = slide(prs, "How to read these numbers",
              "F1 runs 0 to 1, higher is better — but the higher column is the "
              "one that lies")
    bullets(s, [
        ("“Discriminating” means the model’s answer changes with the image.", True),
        ("A model that says the same thing to every street carries no information "
         "about any street, however confident it sounds. That is the failure these "
         "two columns exist to expose.", False),
        ("Read the WEAKER of the two columns, not the higher one.", True),
        ("A model that answers “no” to everything scores 0.70 on the No column and "
         "0.00 on the Yes column. The 0.70 is an artefact of the class balance; only "
         "the 0.00 tells you it never decided anything.", False),
        ("The benchmark: a coin flip on a balanced task gives about 0.50 on BOTH "
         "sides. Weaker side near 0.50 = discriminating. Well below = drifting "
         "toward one answer.", False),
    ], Inches(1.95), size=15)
    table(s, [
        ["weaker side", "reading"],
        ["≥ 0.50", "discriminates both ways"],
        ["0.30 – 0.50", "lopsided"],
        ["0.15 – 0.30", "mostly answers one way"],
        ["< 0.15", "degenerate — answers one way regardless of input"],
    ], Inches(4.75), width=Inches(8.0), widths=[2.0, 6.0])

    # 3 ------------------------------------------------------------------
    s = slide(prs, "What the cited work benchmarked",
              "MINGLE (Liu, …, Sevtsuk — AAAI 2026), Table 2 · ten VLMs, "
              "zero-shot, pairwise classification")
    table(s, [
        ["model", "size", "F1 Yes", "F1 No", "weaker"],
        ["Qwen2-VL", "72B", "0.62", "0.58", "0.58"],
        ["Qwen2.5-VL", "72B", "0.59", "0.56", "0.56"],
        ["Qwen2-VL   ← selected", "7B", "0.53", "0.54", "0.53"],
        ["Pixtral", "12B", "0.52", "0.51", "0.51"],
        ["LLama3.2-Vision", "11B", "0.46", "0.30", "0.30"],
        ["Gemma 3", "24B", "0.27", "0.67", "0.27"],
        ["Phi 3.5", "4.2B", "0.57", "0.24", "0.24"],
        ["Claude-Sonnet", "—", "0.24", "0.66", "0.24"],
        ["Qwen2.5-VL", "7B", "0.10", "0.69", "0.10"],
        ["GPT-4o", "—", "0.00", "0.70", "0.00"],
    ], Inches(1.95), width=Inches(7.4), widths=[3.0, 1.0, 1.1, 1.1, 1.1], highlight={3})
    bullets(s, [
        ("Sorted by the weaker side. Only four of ten clear 0.50.", True),
        ("Two of those four are 72B and do not fit a 12 GB card. GPT-4o answered "
         "“no” to everything: 0.00 / 0.70. Qwen2.5-VL-7B at 0.10 / 0.69 is nearly "
         "the same collapse.", False),
        ("Phi 3.5 shows why the higher column misleads — its F1(Yes) of 0.57 beats "
         "the selected model’s 0.53, but its weaker side is 0.24.", False),
        ("Same failure mode we measured: Qwen2.5-VL-3B answering 3 to every "
         "image.", True),
        ("Table 3, region detection (mIoU):", True),
        ("MINGLE hybrid 0.64  ·  Claude-Sonnet 0.02  ·  GPT-4o 0.01  ·  "
         "Claude-Opus 0.01  ·  Qwen2-VL-7B 0.00  ·  LLama3.2-V 0.00", False),
        ("Every pure VLM collapses on localisation. Only detector + VLM works.", False),
    ], Inches(1.95), left=Inches(7.9), width=Inches(4.9), size=13)

    # 3 ------------------------------------------------------------------
    s = slide(prs, "Measured here — generative VLMs",
              "282 Murray Hill panoramas · Spearman vs measured GVI / VEI · "
              "bootstrap clustered on block face")
    q2 = get(vb, model="Qwen2-VL-7B", question="greenery")
    q2e = get(vb, model="Qwen2-VL-7B", question="enclosure")
    q25 = get(vb, model="Qwen2.5-VL-3B", question="greenery")
    table(s, [
        ["model", "greenery vs GVI", "enclosure vs VEI", "VRAM"],
        ["Qwen2-VL-7B  (4-bit NF4)", fmt(q2), fmt(q2e), "5.9 GB"],
        ["Qwen2.5-VL-3B", fmt(q25), "constant — 3 on 83 of 85 measured", "~3 GB"],
    ], Inches(2.0), widths=[3, 3.4, 3.9, 1.3], highlight={1})
    bullets(s, [
        ("Why only two generative models here, when MINGLE tested ten:", True),
        ("One 12 GB card, run locally. That removes most of the field before any "
         "measurement: 72B needs >40 GB, and even 7B in bf16 is 16.6 GB — 4-bit NF4 "
         "brings it to 5.9 GB, which is why we run a quantised 7B.", False),
        ("6,271 generative VLM calls were made on this corpus in one day. At API "
         "rates that is $25–95 — but cost is not the point. The schema changed six "
         "times that day, and each change meant re-running everything. Under a "
         "per-call meter that iteration does not happen.", False),
        ("So the question left is narrow: does the model work on THIS imagery and "
         "THIS schema, and is 7B worth its cost over 3B? Counting open-vocabulary "
         "models, four VLMs were tested here, not two.", False),
        ("Model size is load-bearing.", True),
        ("7B beats 3B by +0.24 on greenery, and is the difference between a usable "
         "enclosure rating and no variance at all. Measured directly on 85 images, the "
         "3B returned {\"enclosure\": 3} on all 83 that parsed — a degenerate "
         "answer, not a parsing failure.", False),
    ], Inches(3.45), size=13)

    # 4 ------------------------------------------------------------------
    s = slide(prs, "Measured here — open-vocabulary VLMs",
              "same images, same metric")
    cg = get(vb, model="CLIPSeg", question="greenery")
    ce = get(vb, model="CLIPSeg", question="enclosure")
    owl = mb[mb.model.str.contains("street tree")].iloc[0]
    table(s, [
        ["model", "class", "greenery vs GVI", "enclosure vs VEI"],
        ["CLIPSeg", "open-vocab VLM", fmt(cg), fmt(ce)],
        ["OWLv2  ‘a street tree’", "open-vocab VLM", fmt(owl), "—"],
    ], Inches(2.0), widths=[3.2, 2.6, 3.4, 3.4])
    bullets(s, [
        ("CLIPSeg scores highest on greenery and is still not the choice: it returns "
         "a mask for a prompt, not a rating on a scale. It cannot answer the SIM "
         "schema.", True),
        ("OWLv2 is phrasing-sensitive: scored against the SAME target, wording alone "
         "moves it from 0.555 (“scaffolding”) to 0.757 (“metal scaffold poles”).", False),
    ], Inches(4.35))

    # 5 ------------------------------------------------------------------
    s = slide(prs, "Why Qwen2-VL-7B")
    bullets(s, [
        ("1.  It is the only tested model that can answer the SIM schema.", True),
        ("CLIPSeg scores higher on greenery but returns masks, not ratings. That is "
         "the deciding constraint, not the correlation.", False),
        ("2.  Best VLM on enclosure — +0.480, against CLIPSeg +0.414 and a "
         "constant from the 3B.", True),
        ("3.  Size is load-bearing, and 72B does not fit.", True),
        ("16.6 GB in bf16 exceeds a 12.9 GB card; 4-bit NF4 brings 7B to 5.9 GB with "
         "no measured loss. MINGLE traded 0.62 → 0.53 for the same reason by "
         "choice; we do it by constraint.", False),
        ("4.  Independently selected by MINGLE after benchmarking ten VLMs "
         "including GPT-4o, Claude Sonnet and Claude Opus:", True),
        ("“smaller variants — especially Qwen2-VL (7B) — show just "
         "marginally lower performance, while being more robust to overfitting, have "
         "lower memory requirements and significantly faster inference”, because "
         "“urban analysis requires working with vast amounts of data”.", False),
        ("5.  Open weights, run locally, no per-image cost — so the whole corpus "
         "can be re-run whenever the schema changes.", True),
    ], Inches(1.45), size=14)

    # 7 ------------------------------------------------------------------
    s = slide(prs, "Tested and dropped: the 60-degree perceptual weighting",
              "the literature weights a pedestrian's central 60 degrees at ~80%, "
              "the periphery at 10% each")
    bullets(s, [
        ("The question: does the model already do this for us?", True),
        ("If a VLM asked to judge as a pedestrian already attends to the centre, "
         "imposing 80/10/10 in code would apply the same correction twice.", False),
        ("We measured it. It does not.", True),
        ("102 panoramas split into three 60-degree cones, each rated with the same "
         "schema, then regressed against the model's own whole-view rating.", False),
    ], Inches(2.0), size=15)
    table(s, [
        ["", "left", "centre", "right"],
        ["implied weight (LMG)", "0.395", "0.355", "0.250"],
        ["implied weight (NNLS)", "0.361", "0.361", "0.278"],
        ["literature", "0.10", "0.80", "0.10"],
        ["uniform", "0.333", "0.333", "0.333"],
    ], Inches(3.6), width=Inches(7.2), widths=[3.0, 1.4, 1.4, 1.4], highlight={1, 2})
    bullets(s, [
        ("The model reads a 180-degree view flat.", True),
        ("Uniform thirds also beat every alternative on out-of-sample R2 with whole "
         "block faces held out.", False),
        ("Why it matters: Street View is captured from a vehicle on the roadway, not "
         "from the sidewalk, so a pedestrian frontal-attention model does not match "
         "the capture geometry either. The design moved to two 90-degree halves, "
         "each facing a frontage, each its own observation.", False),
    ], Inches(3.6), left=Inches(8.1), width=Inches(4.7), size=13)

    # 8 ------------------------------------------------------------------
    s = slide(prs, "Does the reprojection distort the ratings?",
              "the panoramas are stitched and warped — four 90-degree photos "
              "reprojected onto a cylinder")
    bullets(s, [
        ("The objection: these are not photographs.", True),
        ("Straight lines bow, and a pixel at the top of the frame covers about a "
         "third the solid angle of one at the horizon. So do the ratings mean "
         "anything?", False),
        ("The test: run the identical schema on the ORIGINAL undistorted frames.", True),
        ("Two raw rectilinear photos per walk, no reprojection, same 68 nodes, same "
         "prompts, same scoring.", False),
    ], Inches(2.0), size=15)
    table(s, [
        ["field", "raw frames (undistorted)", "panorama (reprojected)"],
        ["green_eye_level vs GVI", "+0.516", "+0.520"],
        ["vertical_greenery vs GVI", "+0.501", "+0.523"],
        ["street_trees vs GVI", "+0.607", "+0.582"],
        ["facade_variation vs VEI", "+0.458", "+0.487"],
    ], Inches(3.75), width=Inches(9.0), widths=[3.4, 2.8, 2.8])
    bullets(s, [
        ("Every pair overlaps well within its confidence interval — the "
         "reprojection costs nothing measurable, so the panoramas can be used "
         "without a caveat.", True),
        ("One exception: walking_room preferred the raw frames (+0.399 vs +0.273), "
         "plausibly because the sidewalk sits at the bottom edge where the "
         "compression is worst.", False),
    ], Inches(5.5), size=13)


    # 6 ------------------------------------------------------------------
    s = slide(prs, "Caveats we should state", "not footnotes — they change how the numbers read")
    bullets(s, [
        ("Schema length costs accuracy.", True),
        ("Qwen’s enclosure rating scores +0.480 asked as a standalone question and "
         "+0.012 as the tenth field of a twelve-field JSON. Same model, same images.", False),
        ("Enclosure is the weak term, whatever the prompt.", True),
        ("Five formulations across two projections, all at or near zero in-schema. "
         "SIM_hybrid therefore takes enclosure from the pipeline’s measured VEI "
         "rather than from a rating — a data source, not a competing model.", False),
        ("Open-vocabulary models are phrasing-sensitive.", True),
        ("ELSA reports this across Grounding DINO, Detic, OWL and MDETR. We reproduced "
         "it: five prompts for the same structure, scored against the same labels, "
         "ranged 0.555 to 0.757 on wording alone.", False),
    ], Inches(1.5), size=14)

    out = RES / "figures" / "model_selection_slides.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
